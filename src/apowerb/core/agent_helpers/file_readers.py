"""File reading helpers for images, archives, Office documents and binaries."""
from __future__ import annotations

import os

from apowerb.core.agent_helpers.text_utils import _truncate_content


_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".ico", ".webp"}
_AUDIO_EXTENSIONS = {".mp3", ".wav", ".ogg", ".m4a", ".webm", ".flac", ".aac", ".wma", ".opus"}
_ARCHIVE_EXTENSIONS = {".zip"}
_MAX_IMAGE_BASE64_SIZE = 200 * 1024  # base64 encode images up to 200 KB


def _is_binary_by_magic(file_path: str) -> bool:
    """Check if a file is binary by looking at the first bytes for null characters."""
    try:
        with open(file_path, "rb") as f:
            chunk = f.read(8192)
        return b"\x00" in chunk
    except OSError:
        return False


def _read_image(file_path: str, file_size: int, filename: str) -> dict:
    """Return image metadata and base64 content (if small enough)."""
    import base64

    result = {
        "status": "success",
        "filename": filename,
        "size": file_size,
        "format": "image",
        "content_type": "image/"
        + os.path.splitext(filename)[1].lstrip(".").lower().replace("jpg", "jpeg"),
    }
    if file_size <= _MAX_IMAGE_BASE64_SIZE:
        with open(file_path, "rb") as f:
            raw = f.read()
        b64 = base64.b64encode(raw).decode("ascii")
        result["base64"] = b64
        result["message"] = (
            "Image encoded as base64. You can describe it or reference it."
        )
    else:
        result["message"] = (
            f"Image too large for inline base64 ({file_size} bytes). "
            "File is available on disk. Describe what you know or ask the user for details."
        )
    return result


def _read_archive(file_path: str, file_size: int, filename: str) -> dict:
    """List contents of a ZIP archive."""
    import zipfile

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            entries = []
            for info in zf.infolist():
                entries.append(f"  {info.filename}  ({info.file_size} bytes)")
            listing = "\n".join(entries[:200])  # cap at 200 entries
            total = len(zf.infolist())
            if total > 200:
                listing += f"\n  ... and {total - 200} more files"
        return {
            "status": "success",
            "filename": filename,
            "size": file_size,
            "format": "zip_listing",
            "content": f"ZIP archive — {total} file(s):\n{listing}",
            "total_files": total,
        }
    except Exception as e:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": f"Could not read ZIP archive: {e}",
        }


def _read_excel(file_path: str, file_size: int, filename: str) -> dict:
    """Extract content from .xlsx / .xls files using openpyxl or xlrd."""
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets_content = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    # Skip completely empty rows
                    if any(cell is not None for cell in row):
                        rows.append("\t".join("" if v is None else str(v) for v in row))
                if rows:
                    sheets_content.append(f"### Sheet: {sheet_name}\n" + "\n".join(rows))
            wb.close()
            content = "\n\n".join(sheets_content) or "(No data found in workbook)"
        else:
            # .xls — use xlrd
            import xlrd
            wb = xlrd.open_workbook(file_path)
            sheets_content = []
            for sheet in wb.sheets():
                rows = []
                for r in range(sheet.nrows):
                    row_vals = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                    rows.append("\t".join(row_vals))
                if rows:
                    sheets_content.append(f"### Sheet: {sheet.name}\n" + "\n".join(rows))
            content = "\n\n".join(sheets_content) or "(No data found in workbook)"

        content, truncated = _truncate_content(content)
        return {
            "status": "success",
            "filename": filename,
            "size": file_size,
            "format": "excel",
            "content": content,
            "truncated": truncated,
        }
    except ImportError as e:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": f"Missing library to read Excel files: {e}. Install openpyxl (xlsx) or xlrd (xls).",
        }
    except Exception as e:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": f"Could not read Excel file: {e}",
        }


def _read_docx(file_path: str, file_size: int, filename: str) -> dict:
    """Extract text from a .docx Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        content = "\n\n".join(paragraphs) or "(No text content found)"
        content, truncated = _truncate_content(content)
        return {
            "status": "success",
            "filename": filename,
            "size": file_size,
            "format": "docx",
            "content": content,
            "truncated": truncated,
        }
    except ImportError:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": "Missing library to read Word files. Install python-docx.",
        }
    except Exception as e:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": f"Could not read Word file: {e}",
        }


def _read_pptx(file_path: str, file_size: int, filename: str) -> dict:
    """Extract text from a .pptx PowerPoint presentation."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_content.append(f"### Slide {i}\n" + "\n".join(texts))
        content = "\n\n".join(slides_content) or "(No text content found)"
        content, truncated = _truncate_content(content)
        return {
            "status": "success",
            "filename": filename,
            "size": file_size,
            "format": "pptx",
            "content": content,
            "truncated": truncated,
        }
    except ImportError:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": "Missing library to read PowerPoint files. Install python-pptx.",
        }
    except Exception as e:
        return {
            "status": "error",
            "filename": filename,
            "size": file_size,
            "message": f"Could not read PowerPoint file: {e}",
        }


def _read_binary_metadata(file_path: str, file_size: int, filename: str) -> dict:
    """Return basic metadata for unknown binary files."""
    ext = os.path.splitext(filename)[1].lower()
    return {
        "status": "success",
        "filename": filename,
        "size": file_size,
        "format": "binary_metadata",
        "message": (
            f"Binary file ({ext or 'unknown type'}, {file_size} bytes). "
            "Cannot extract text content. The file is stored and available on disk."
        ),
    }
